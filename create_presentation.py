import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()
    
    # Set slide dimensions (16:9 Widescreen)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Colors
    DARK_BG = RGBColor(15, 23, 42)      # #0F172A
    WHITE = RGBColor(255, 255, 255)      # #FFFFFF
    SKY_BLUE = RGBColor(56, 189, 248)    # #38BDF8
    LIGHT_GRAY = RGBColor(148, 163, 184) # #94A6B8
    GREEN_ACCENT = RGBColor(16, 185, 129) # #10B981
    
    # Helper to set slide background to Dark Navy
    def set_slide_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = DARK_BG

    # Helper to add standard text box
    def add_textbox(slide, left, top, width, height):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        return tf

    blank_slide_layout = prs.slide_layouts[6] # Blank layout

    # ==========================================
    # SLIDE 1: Title Slide (Cover)
    # ==========================================
    slide1 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide1)
    
    tf1 = add_textbox(slide1, Inches(1), Inches(1.8), Inches(11.3), Inches(4.5))
    
    p_title = tf1.paragraphs[0]
    p_title.text = "SIMPLE LMS EXTENDED BACKEND"
    p_title.font.name = "Arial"
    p_title.font.size = Pt(44)
    p_title.font.bold = True
    p_title.font.color.rgb = WHITE
    p_title.alignment = PP_ALIGN.LEFT
    
    p_subtitle = tf1.add_paragraph()
    p_subtitle.text = "Pengembangan Backend Teraugmentasi & Integrasi REST API Django Ninja"
    p_subtitle.font.name = "Arial"
    p_subtitle.font.size = Pt(22)
    p_subtitle.font.color.rgb = SKY_BLUE
    p_subtitle.space_before = Pt(15)
    
    p_info = tf1.add_paragraph()
    p_info.text = "\nMata Kuliah: Pemrograman Sisi Server (A11.54403)\nNama: Adi Priyo  |  NIM: [Lengkapi NIM Anda]\nUniversitas Dian Nuswantoro"
    p_info.font.name = "Arial"
    p_info.font.size = Pt(16)
    p_info.font.color.rgb = LIGHT_GRAY
    p_info.space_before = Pt(40)

    # ==========================================
    # SLIDE 2: Arsitektur & Teknologi
    # ==========================================
    slide2 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide2)
    
    tf2_title = add_textbox(slide2, Inches(1), Inches(0.6), Inches(11.3), Inches(1))
    p = tf2_title.paragraphs[0]
    p.text = "Arsitektur & Desain Sistem"
    p.font.name = "Arial"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    tf2_body = add_textbox(slide2, Inches(1), Inches(1.8), Inches(11.3), Inches(5))
    
    bullet1 = tf2_body.paragraphs[0]
    bullet1.text = "• Model Pembagian Kontainer (Docker Compose):"
    bullet1.font.name = "Arial"
    bullet1.font.size = Pt(20)
    bullet1.font.bold = True
    bullet1.font.color.rgb = SKY_BLUE
    
    sub1_1 = tf2_body.add_paragraph()
    sub1_1.text = "  - Service Web: Django Application (Django Ninja & Python 3.11)\n  - Service Database: PostgreSQL 15 Container\n  - Service Caching Layer: Redis Container"
    sub1_1.font.name = "Arial"
    sub1_1.font.size = Pt(16)
    sub1_1.font.color.rgb = LIGHT_GRAY
    sub1_1.space_after = Pt(15)

    bullet2 = tf2_body.add_paragraph()
    bullet2.text = "• Pola Integrasi REST API & Caching:"
    bullet2.font.name = "Arial"
    bullet2.font.size = Pt(20)
    bullet2.font.bold = True
    bullet2.font.color.rgb = SKY_BLUE
    
    sub2_1 = tf2_body.add_paragraph()
    sub2_1.text = "  - Request Client divalidasi otomatis menggunakan Pydantic Schema.\n  - JWT Bearer Token mengamankan komunikasi data (Stateless Auth).\n  - Redis bertindak sebagai memori cepat untuk menyimpan list/detail course."
    sub2_1.font.name = "Arial"
    sub2_1.font.size = Pt(16)
    sub2_1.font.color.rgb = LIGHT_GRAY

    # ==========================================
    # SLIDE 3: Kriteria Fondasi Wajib (30 Poin)
    # ==========================================
    slide3 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide3)
    
    tf3_title = add_textbox(slide3, Inches(1), Inches(0.6), Inches(11.3), Inches(1))
    p = tf3_title.paragraphs[0]
    p.text = "Pembuktian Fondasi Wajib (30 Poin)"
    p.font.name = "Arial"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    tf3_body = add_textbox(slide3, Inches(1), Inches(1.8), Inches(11.3), Inches(5))
    
    bullets = [
        ("✔ Docker & Docker Compose", "Peletakan: docker-compose.yml & Dockerfile di root project.", SKY_BLUE),
        ("✔ Database PostgreSQL & Migrasi", "Peletakan: settings.py (Logika koneksi DB & pooling).", SKY_BLUE),
        ("✔ JWT Authentication & Authorization", "Peletakan: lms/auth.py (JWT Bearer Token handler & class JWTAuth).", SKY_BLUE),
        ("✔ Role-Based Access Control (RBAC)", "Peletakan: lms/models.py (Field role pada class User) & lms/api.py (Validasi role).", SKY_BLUE),
        ("✔ Endpoint CRUD Dasar", "Peletakan: lms/api.py (Routers for Courses, Lessons, Enrollments, Progress).", SKY_BLUE),
        ("✔ Swagger/OpenAPI Docs & .env Config", "Peletakan: /api/docs & .env.example (Variabel rahasia dipisah).", SKY_BLUE)
    ]
    
    for i, (title, desc, color) in enumerate(bullets):
        p_b = tf3_body.paragraphs[0] if i == 0 else tf3_body.add_paragraph()
        p_b.text = f"{title}"
        p_b.font.name = "Arial"
        p_b.font.size = Pt(18)
        p_b.font.bold = True
        p_b.font.color.rgb = color
        if i > 0:
            p_b.space_before = Pt(10)
            
        p_d = tf3_body.add_paragraph()
        p_d.text = f"  -> {desc}"
        p_d.font.name = "Arial"
        p_d.font.size = Pt(14)
        p_d.font.color.rgb = LIGHT_GRAY

    # ==========================================
    # SLIDE 4: Kriteria Fitur Tambahan (46 Poin)
    # ==========================================
    slide4 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide4)
    
    tf4_title = add_textbox(slide4, Inches(1), Inches(0.6), Inches(11.3), Inches(1))
    p = tf4_title.paragraphs[0]
    p.text = "Fitur Tambahan Terpilih (46 Poin)"
    p.font.name = "Arial"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    tf4_body = add_textbox(slide4, Inches(1), Inches(1.8), Inches(11.3), Inches(5))
    
    f_add = [
        ("1. Permission & Ownership Ketat (12 Poin)", "Peletakan: lms/api.py\n- Memvalidasi agar Instruktur hanya bisa mengubah/menghapus course miliknya.\n- Membatasi agar Siswa harus terdaftar (enrolled) untuk melihat materi lesson."),
        ("2. Search, Filter, & Sorting Lanjutan + Redis Caching (12 Poin)", "Peletakan: lms/api.py\n- Endpoint list_courses memproses parameter query search, filter kategori, & sort (rating/popular).\n- Caching diimplementasikan pada Redis dengan invalidasi otomatis saat data course diperbarui."),
        ("3. Rating, Review, & Wishlist (12 Poin)", "Peletakan: lms/models.py (Review & Wishlist models) & lms/api.py\n- Menyediakan endpoint bagi enrolled student untuk memberi rating bintang & ulasan, dan simpan wishlist."),
        ("4. Course Announcement (10 Poin)", "Peletakan: lms/models.py (Announcement model) & lms/api.py\n- Endpoint bagi Instruktur untuk membuat pengumuman satu arah khusus untuk para siswa terdaftar.")
    ]
    
    for i, (title, desc) in enumerate(f_add):
        p_b = tf4_body.paragraphs[0] if i == 0 else tf4_body.add_paragraph()
        p_b.text = title
        p_b.font.name = "Arial"
        p_b.font.size = Pt(18)
        p_b.font.bold = True
        p_b.font.color.rgb = GREEN_ACCENT
        if i > 0:
            p_b.space_before = Pt(12)
            
        p_d = tf4_body.add_paragraph()
        p_d.text = desc
        p_d.font.name = "Arial"
        p_d.font.size = Pt(13)
        p_d.font.color.rgb = LIGHT_GRAY

    # ==========================================
    # SLIDE 5: Demonstrasi UI & Bonus (+5 Poin)
    # ==========================================
    slide5 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide5)
    
    tf5_title = add_textbox(slide5, Inches(1), Inches(0.6), Inches(11.3), Inches(1))
    p = tf5_title.paragraphs[0]
    p.text = "Dashboard UI & Bonus Poin (+5 Poin)"
    p.font.name = "Arial"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    tf5_body = add_textbox(slide5, Inches(1), Inches(1.8), Inches(11.3), Inches(5))
    
    b1 = tf5_body.paragraphs[0]
    b1.text = "• Halaman Frontend Dashboard Interaktif (index.html):"
    b1.font.name = "Arial"
    b1.font.size = Pt(20)
    b1.font.bold = True
    b1.font.color.rgb = SKY_BLUE
    
    b1_desc = tf5_body.add_paragraph()
    b1_desc.text = "  - Lokasi: index.html di root project.\n  - Berupa Single Page Application (SPA) minimalis untuk visualisasi presentasi.\n  - Mendukung login manual menggunakan akun demo yang tersedia.\n  - Menampilkan fitur RBAC dinamis: Tombol CRUD edit/hapus hanya muncul pada peran Admin/Instructor."
    b1_desc.font.name = "Arial"
    b1_desc.font.size = Pt(16)
    b1_desc.font.color.rgb = LIGHT_GRAY
    b1_desc.space_after = Pt(15)

    b2 = tf5_body.add_paragraph()
    b2.text = "• Integrasi Postman Collection:"
    b2.font.name = "Arial"
    b2.font.size = Pt(20)
    b2.font.bold = True
    b2.font.color.rgb = SKY_BLUE
    
    b2_desc = tf5_body.add_paragraph()
    b2_desc.text = "  - Lokasi: LMS_Postman_Collection.json di root project.\n  - Memudahkan pengujian manual seluruh API end-to-end bagi tim penguji/dosen."
    b2_desc.font.name = "Arial"
    b2_desc.font.size = Pt(16)
    b2_desc.font.color.rgb = LIGHT_GRAY

    # ==========================================
    # SLIDE 6: Panduan Menjalankan Proyek (Cara Running)
    # ==========================================
    slide6 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide6)
    
    tf6_title = add_textbox(slide6, Inches(1), Inches(0.6), Inches(11.3), Inches(1))
    p = tf6_title.paragraphs[0]
    p.text = "Panduan Menjalankan Proyek"
    p.font.name = "Arial"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    tf6_body = add_textbox(slide6, Inches(1), Inches(1.8), Inches(11.3), Inches(5))
    
    c1 = tf6_body.paragraphs[0]
    c1.text = "A. Menggunakan Docker Compose (Kontainerisasi)"
    c1.font.name = "Arial"
    c1.font.size = Pt(18)
    c1.font.bold = True
    c1.font.color.rgb = SKY_BLUE
    
    c1_desc = tf6_body.add_paragraph()
    c1_desc.text = "   1. cp .env.example .env\n   2. docker compose up --build -d\n   3. docker compose exec web python manage.py migrate\n   4. docker compose exec web python manage.py seed_data"
    c1_desc.font.name = "Consolas"
    c1_desc.font.size = Pt(13)
    c1_desc.font.color.rgb = LIGHT_GRAY
    c1_desc.space_after = Pt(15)

    c2 = tf6_body.add_paragraph()
    c2.text = "B. Menggunakan Environment Lokal (SQLite Fallback)"
    c2.font.name = "Arial"
    c2.font.size = Pt(18)
    c2.font.bold = True
    c2.font.color.rgb = SKY_BLUE
    
    c2_desc = tf6_body.add_paragraph()
    c2_desc.text = "   1. .\\venv\\Scripts\\activate\n   2. python manage.py migrate\n   3. python manage.py seed_data\n   4. python manage.py runserver\n   5. python demo_simulation.py (untuk menjalankan demo API interaktif)"
    c2_desc.font.name = "Consolas"
    c2_desc.font.size = Pt(13)
    c2_desc.font.color.rgb = LIGHT_GRAY

    # Save presentation
    output_path = r"c:\Users\Adi Priyo\Documents\final projek pss\Presentation1_Extended.pptx"
    try:
        prs.save(output_path)
        print(f"Presentation successfully created and saved to {output_path}")
    except Exception as e:
        print(f"Error saving file: {e}")

if __name__ == "__main__":
    create_presentation()
